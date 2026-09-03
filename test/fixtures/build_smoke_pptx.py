"""Build a small PPTX for browser smoke testing of extract / inject."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUTPUTS = [
    ROOT / "test" / "fixtures" / "smoke-test.pptx",
]


def set_run(paragraph, text, *, bold=False, color=None, size=18):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return run


def add_textbox(slide, left, top, width, height, builder, *, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    builder(p, font_size)
    return box


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    add_textbox(
        slide1, 0.6, 0.35, 12, 0.7,
        lambda p, size: set_run(p, "Sample Deck: Observability Sales Play", bold=True, size=28),
    )
    add_textbox(
        slide1, 0.6, 1.2, 12, 0.5,
        lambda p, size: set_run(p, "Who to talk to", bold=True, size=20),
    )
    add_textbox(
        slide1, 0.6, 1.8, 12, 0.8,
        lambda p, size: (
            set_run(p, "AI teams who ", size=18),
            set_run(p, "build, operate and improve", bold=True, color=RGBColor(0x1D, 0x4E, 0xD8), size=18),
            set_run(p, " AI agents & LLMs", size=18),
        ),
    )
    add_textbox(
        slide1, 0.6, 2.8, 5.8, 0.4,
        lambda p, size: (
            set_run(p, "Network ", color=RGBColor(0x16, 0xA3, 0x4A), bold=True, size=22),
            set_run(p, "Observability", color=RGBColor(0x15, 0x80, 0x3D), bold=True, size=22),
        ),
    )
    add_textbox(
        slide1, 6.8, 2.8, 5.8, 0.4,
        lambda p, size: (
            set_run(p, "Full-Stack ", color=RGBColor(0x25, 0x63, 0xEB), bold=True, size=22),
            set_run(p, "Observability", color=RGBColor(0x1D, 0x4E, 0xD8), bold=True, size=22),
        ),
    )

    table_shape = slide1.shapes.add_table(2, 2, Inches(0.6), Inches(3.6), Inches(7.2), Inches(1.6))
    table = table_shape.table
    table.cell(0, 0).text = "Role"
    table.cell(0, 1).text = "Person"
    table.cell(1, 0).text = "Economic Buyer"
    table.cell(1, 1).text = "CTO / Head of AI Eng."
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)

    add_textbox(
        slide1, 8.1, 3.6, 4.6, 1.6,
        lambda p, size: set_run(
            p,
            "Build trust in AI by evaluating agent behavior and token costs.",
            size=16,
        ),
    )
    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = "Speaker note: keep product names such as Splunk and ThousandEyes."

    slide2 = prs.slides.add_slide(blank)
    add_textbox(
        slide2, 0.6, 0.4, 12, 0.6,
        lambda p, size: set_run(p, "How our sales plays align to your customer", bold=True, size=26),
    )
    add_textbox(
        slide2, 0.6, 1.3, 12, 0.8,
        lambda p, size: set_run(
            p,
            "Expand by co-selling Observability Cloud and ThousandEyes.",
            size=18,
        ),
    )
    add_textbox(
        slide2, 0.6, 2.3, 3.8, 0.7,
        lambda p, size: (
            set_run(p, "Capture", bold=True, color=RGBColor(0xB4, 0x53, 0x09), size=24),
        ),
    )
    add_textbox(
        slide2, 0.6, 3.0, 3.8, 0.5,
        lambda p, size: set_run(p, "new market", size=18),
    )
    add_textbox(
        slide2, 4.7, 2.3, 3.8, 0.7,
        lambda p, size: set_run(p, "Grow", bold=True, color=RGBColor(0x1D, 0x4E, 0xD8), size=24),
    )
    add_textbox(
        slide2, 4.7, 3.0, 3.8, 0.5,
        lambda p, size: set_run(p, "the business", size=18),
    )
    add_textbox(
        slide2, 8.8, 2.3, 3.8, 0.7,
        lambda p, size: set_run(p, "Protect", bold=True, color=RGBColor(0x16, 0x65, 0x34), size=24),
    )
    add_textbox(
        slide2, 8.8, 3.0, 3.8, 0.5,
        lambda p, size: set_run(p, "the base", size=18),
    )
    add_textbox(
        slide2, 0.6, 4.2, 12, 1.4,
        lambda p, size: set_run(
            p,
            "Empty notes on this slide are intentional. Tag boundaries like [0]...[/0] must survive translation.",
            size=16,
        ),
    )

    for path in OUTPUTS:
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        print(path)


if __name__ == "__main__":
    build()
